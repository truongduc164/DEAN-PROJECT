"""
tests/test_ppt_processor.py – Tests for PPT translation processor.
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from pptx import Presentation
from pptx.util import Inches

from app.core.ppt_processor import PptProcessor, build_ppt_output_path
from app.core.translators.translator_service import MockTranslator
from app.core.ocr.models import ImageOcrResult, OcrSegment


class TestBuildPptOutputPath:
    def test_vietnamese_suffix(self, tmp_path: Path):
        p = tmp_path / "demo.pptx"
        out = build_ppt_output_path(p, "Vietnamese")
        assert out.name == "demo_Vi.pptx"

    def test_english_suffix(self, tmp_path: Path):
        p = tmp_path / "demo.pptx"
        out = build_ppt_output_path(p, "English")
        assert out.name == "demo_En.pptx"


class TestPptProcessor:
    def _make_pptx(self, path: Path, texts: list[str]) -> Path:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        for i, t in enumerate(texts):
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(0.5))
            tf = txBox.text_frame
            tf.text = t
        prs.save(str(path))
        return path

    def test_basic_translation(self, tmp_path: Path):
        src = self._make_pptx(tmp_path / "test.pptx", ["Hello world", "Good morning"])
        proc = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
        )
        out = proc.process(src)
        assert out.exists()
        assert "_Vi" in out.name

        prs_out = Presentation(str(out))
        slide = prs_out.slides[0]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        # MockTranslator changes text
        assert any(t != "Hello world" for t in texts)

    def test_prefix_mode(self, tmp_path: Path):
        src = self._make_pptx(tmp_path / "prefix.pptx", ["Hello"])
        proc = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
            output_mode="prefix",
        )
        out = proc.process(src)
        prs_out = Presentation(str(out))
        for shape in prs_out.slides[0].shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Hello" in text:
                    assert text.endswith("Hello")

    def test_suffix_mode(self, tmp_path: Path):
        src = self._make_pptx(tmp_path / "suffix.pptx", ["Hello"])
        proc = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
            output_mode="suffix",
        )
        out = proc.process(src)
        prs_out = Presentation(str(out))
        for shape in prs_out.slides[0].shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Hello" in text:
                    assert text.startswith("Hello")

    def test_skip_textboxes(self, tmp_path: Path):
        """When translate_textboxes=False, textboxes should be skipped."""
        src = self._make_pptx(tmp_path / "notb.pptx", ["Keep this"])
        proc = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
            translate_textboxes=False,
        )
        out = proc.process(src)
        prs_out = Presentation(str(out))
        # Textboxes should NOT be translated (no placeholder_format)
        for shape in prs_out.slides[0].shapes:
            if getattr(shape, "has_text_frame", False):
                try:
                    _ = shape.placeholder_format
                    is_placeholder = True
                except (ValueError, AttributeError):
                    is_placeholder = False
                
                if not is_placeholder:
                    # Text should remain unchanged
                    assert "Keep this" in shape.text_frame.text

    def test_cancel_creates_stop_file(self, tmp_path: Path):
        import threading
        cancel = threading.Event()

        class CancellingTranslator(MockTranslator):
            def translate(self, text, src, tgt):
                cancel.set()  # Cancel during translation
                return super().translate(text, src, tgt)

        src = self._make_pptx(tmp_path / "cancel.pptx", ["Text", "More text"])
        proc = PptProcessor(
            translator=CancellingTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
            cancel_event=cancel,
        )
        out = proc.process(src)
        assert "_stop" in out.name

    def test_empty_presentation(self, tmp_path: Path):
        src = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(str(src))

        proc = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
        )
        out = proc.process(src)
        assert out.exists()


class _DummyEventManager:
    def __init__(self):
        self.logs: list[tuple[str, str]] = []

    def log(self, level: str, msg: str):
        self.logs.append((level, msg))

    def progress(self, *_args, **_kwargs):
        pass


class TestPptOcrModes:
    def _make_proc(self) -> tuple[PptProcessor, _DummyEventManager]:
        em = _DummyEventManager()
        proc = PptProcessor(
            translator=MockTranslator(),
            event_manager=em,
            source_lang="Chinese",
            target_lang="Vietnamese",
        )
        return proc, em

    def test_mode_pairs_are_independent(self):
        proc, _ = self._make_proc()
        pairs = [
            ("custom_format", "exact"),
            ("custom_format", "smart_adjust"),
            ("custom_format", "whitespace"),
            ("keep_original_format", "whitespace"),
        ]
        for fmt, place in pairs:
            assert proc._resolve_text_format_mode(fmt, True) == fmt
            assert proc._resolve_ocr_placement_mode(place, None) == place

    def test_whitespace_success_does_not_trigger_smart_adjust(self, monkeypatch):
        proc, em = self._make_proc()
        base = (100, 100, 80, 30)
        occ = [(0, 0, 90, 90)]
        shape_rects = []
        source_rect = (50, 50, 120, 80)

        monkeypatch.setattr(proc, "_find_same_section_region", lambda *_args, **_kwargs: None)

        monkeypatch.setattr(
            proc,
            "_choose_whitespace_position",
            lambda *_args, **_kwargs: ((120, 100, 80, 30), "right"),
        )

        def _must_not_call(*_args, **_kwargs):
            raise AssertionError("smart_adjust should not run when whitespace succeeds")

        monkeypatch.setattr(proc, "_choose_smart_position", _must_not_call)

        rect, overlap = proc._resolve_placement_rect(
            base,
            occ,
            shape_rects,
            source_rect,
            "whitespace",
            step_x=10,
            step_y=10,
            max_shift_x=80,
            max_shift_y=80,
            whitespace_search_x=180,
            whitespace_search_y=180,
        )
        assert rect == (120, 100, 80, 30)
        assert overlap is False
        assert any("placement=nearby_whitespace" in msg for _lvl, msg in em.logs)

    def test_same_section_whitespace_preferred_over_nearby(self, monkeypatch):
        proc, em = self._make_proc()
        base = (100, 100, 80, 30)
        occ = [(0, 0, 90, 90)]
        shape_rects = [(50, 70, 300, 160)]
        source_rect = (60, 90, 120, 80)
        calls = {"nearby": 0}

        monkeypatch.setattr(proc, "_find_same_section_region", lambda *_args, **_kwargs: (50, 70, 300, 160))

        def _ws(*_args, **kwargs):
            if kwargs.get("constrain_region") is not None:
                return (130, 100, 80, 30), "right"
            calls["nearby"] += 1
            return (200, 200, 80, 30), "down"

        monkeypatch.setattr(proc, "_choose_whitespace_position", _ws)

        def _must_not_call(*_args, **_kwargs):
            raise AssertionError("smart_adjust should not run when same-section whitespace succeeds")

        monkeypatch.setattr(proc, "_choose_smart_position", _must_not_call)

        rect, overlap = proc._resolve_placement_rect(
            base,
            occ,
            shape_rects,
            source_rect,
            "whitespace",
            step_x=10,
            step_y=10,
            max_shift_x=80,
            max_shift_y=80,
            whitespace_search_x=180,
            whitespace_search_y=180,
        )
        assert rect == (130, 100, 80, 30)
        assert overlap is False
        assert calls["nearby"] == 0
        assert any("placement=same_section_whitespace" in msg for _lvl, msg in em.logs)

    def test_whitespace_failure_falls_back_to_smart_adjust(self, monkeypatch):
        proc, em = self._make_proc()
        base = (100, 100, 80, 30)
        occ = [(90, 90, 120, 60)]
        shape_rects = []
        source_rect = (50, 50, 120, 80)

        monkeypatch.setattr(proc, "_find_same_section_region", lambda *_args, **_kwargs: None)

        monkeypatch.setattr(
            proc,
            "_choose_whitespace_position",
            lambda *_args, **_kwargs: (None, None),
        )
        called = {"smart": False}

        def _smart(*_args, **_kwargs):
            called["smart"] = True
            return (130, 110, 80, 30), False

        monkeypatch.setattr(proc, "_choose_smart_position", _smart)

        rect, overlap = proc._resolve_placement_rect(
            base,
            occ,
            shape_rects,
            source_rect,
            "whitespace",
            step_x=10,
            step_y=10,
            max_shift_x=80,
            max_shift_y=80,
            whitespace_search_x=180,
            whitespace_search_y=180,
        )
        assert called["smart"] is True
        assert rect == (130, 110, 80, 30)
        assert overlap is False
        assert any("fallback=smart_adjust" in msg for _lvl, msg in em.logs)

    def test_smart_adjust_failure_falls_back_to_exact(self, monkeypatch):
        proc, em = self._make_proc()
        base = (100, 100, 80, 30)
        occ = [(90, 90, 120, 60)]
        shape_rects = []
        source_rect = (50, 50, 120, 80)

        monkeypatch.setattr(
            proc,
            "_choose_smart_position",
            lambda *_args, **_kwargs: ((130, 110, 80, 30), True),
        )

        rect, overlap = proc._resolve_placement_rect(
            base,
            occ,
            shape_rects,
            source_rect,
            "smart_adjust",
            step_x=10,
            step_y=10,
            max_shift_x=80,
            max_shift_y=80,
            whitespace_search_x=180,
            whitespace_search_y=180,
        )
        assert rect == base
        assert overlap is False
        assert any("fallback=exact" in msg for _lvl, msg in em.logs)

    def test_logs_show_text_and_placement_modes(self, monkeypatch):
        proc, em = self._make_proc()
        from app.settings.settings_manager import settings

        cfg = {
            "ocr_settings.render_textbox_overlay": True,
            "ocr_settings.ocr_display_mode": "overwrite",
            "ocr_settings.ocr_display_container": "textbox",
            "ocr_settings.ocr_textbox_placement_mode": "exact",
            "text_style_settings.translated_text_format_mode": "custom_format",
            "text_style_settings.font_family": "Times New Roman",
            "text_style_settings.font_size": 12,
            "text_style_settings.font_color": "#FF0000",
            "text_style_settings.bold": True,
            "text_style_settings.italic": False,
            "text_style_settings.underline": True,
            "ocr_settings.smart_adjust_max_shift_px": 80,
            "ocr_settings.whitespace_search_max_shift_px": 180,
        }
        monkeypatch.setattr(settings, "get", lambda key, default=None: cfg.get(key, default))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        fake_result = ImageOcrResult(
            image_id="s1_1",
            width=100,
            height=100,
            blocks=[
                OcrSegment(
                    id="b1",
                    text="中文",
                    translated_text="[Vietnamese] 中文",
                    confidence=99.0,
                    bbox=(10, 10, 40, 20),
                )
            ],
        )

        monkeypatch.setattr(
            proc.ocr_pipeline,
            "ocr_and_filter_images",
            lambda *_args, **_kwargs: [{"slide": slide, "shape": shape, "result": fake_result, "translatable_blocks": []}],
        )
        monkeypatch.setattr(
            proc.ocr_pipeline,
            "batch_translate_segments",
            lambda *_args, **_kwargs: None,
        )

        proc._process_images([(slide, shape, b"ignored", 1)])
        assert any("text_format_mode=custom_format" in msg for _lvl, msg in em.logs)
        assert any("placement_mode=exact" in msg for _lvl, msg in em.logs)
