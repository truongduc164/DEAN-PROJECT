"""
tests/test_output_naming.py – Tests for the shared output naming helper.

Covers:
    - Basic suffix generation (Vi, En, Ja, …)
    - Collision avoidance: same source → _Vi, _Vi(1), _Vi(2)
    - Works across file types: .pptx, .xlsx, .docx
    - output_dir override
    - PPT processor integration: two translations never overwrite
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from app.core.output_naming import (
    build_output_path,
    get_unique_output_path,
    get_lang_suffix,
    LANG_SUFFIX,
)


# ═══════════════════════════════════════════════════════════════════════
# Unit: get_lang_suffix
# ═══════════════════════════════════════════════════════════════════════

class TestGetLangSuffix:
    def test_known_languages(self):
        assert get_lang_suffix("Vietnamese") == "Vi"
        assert get_lang_suffix("English") == "En"
        assert get_lang_suffix("Japanese") == "Ja"
        assert get_lang_suffix("Chinese") == "Zh"
        assert get_lang_suffix("Korean") == "Ko"

    def test_unknown_language_uses_first_two_chars(self):
        assert get_lang_suffix("French") == "Fr"
        assert get_lang_suffix("Spanish") == "Sp"


# ═══════════════════════════════════════════════════════════════════════
# Unit: build_output_path
# ═══════════════════════════════════════════════════════════════════════

class TestBuildOutputPath:
    def test_vietnamese_pptx(self, tmp_path: Path):
        p = tmp_path / "demo.pptx"
        out = build_output_path(p, "Vietnamese")
        assert out.name == "demo_Vi.pptx"
        assert out.parent == tmp_path

    def test_english_xlsx(self, tmp_path: Path):
        p = tmp_path / "report.xlsx"
        out = build_output_path(p, "English")
        assert out.name == "report_En.xlsx"

    def test_japanese_docx(self, tmp_path: Path):
        p = tmp_path / "letter.docx"
        out = build_output_path(p, "Japanese")
        assert out.name == "letter_Ja.docx"

    def test_custom_output_dir(self, tmp_path: Path):
        p = tmp_path / "source" / "file.pptx"
        out_dir = tmp_path / "output"
        out = build_output_path(p, "Vietnamese", output_dir=out_dir)
        assert out.parent == out_dir
        assert out.name == "file_Vi.pptx"

    def test_extension_override(self, tmp_path: Path):
        p = tmp_path / "file.pptx"
        out = build_output_path(p, "Vietnamese", extension=".pdf")
        assert out.name == "file_Vi.pdf"


# ═══════════════════════════════════════════════════════════════════════
# Unit: get_unique_output_path
# ═══════════════════════════════════════════════════════════════════════

class TestGetUniqueOutputPath:
    def test_first_time_no_collision(self, tmp_path: Path):
        """First translation → input_Vi.pptx (no existing file)."""
        src = tmp_path / "input.pptx"
        src.touch()

        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "input_Vi.pptx"
        assert not out.exists()

    def test_second_time_creates_numbered(self, tmp_path: Path):
        """Second translation → input_Vi(1).pptx."""
        src = tmp_path / "input.pptx"
        src.touch()

        # Simulate first translation output exists
        first = tmp_path / "input_Vi.pptx"
        first.touch()

        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "input_Vi(1).pptx"

    def test_third_time_creates_numbered_2(self, tmp_path: Path):
        """Third translation → input_Vi(2).pptx."""
        src = tmp_path / "input.pptx"
        src.touch()

        (tmp_path / "input_Vi.pptx").touch()
        (tmp_path / "input_Vi(1).pptx").touch()

        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "input_Vi(2).pptx"

    def test_many_collisions(self, tmp_path: Path):
        """Handles many existing files."""
        src = tmp_path / "file.pptx"
        src.touch()

        (tmp_path / "file_Vi.pptx").touch()
        for i in range(1, 6):
            (tmp_path / f"file_Vi({i}).pptx").touch()

        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "file_Vi(6).pptx"

    def test_english_collision(self, tmp_path: Path):
        """Collision works for _En suffix too."""
        src = tmp_path / "data.xlsx"
        src.touch()

        (tmp_path / "data_En.xlsx").touch()

        out = get_unique_output_path(src, "English")
        assert out.name == "data_En(1).xlsx"

    def test_preserves_existing_files(self, tmp_path: Path):
        """Existing translated files are never touched."""
        src = tmp_path / "doc.pptx"
        src.touch()

        first = tmp_path / "doc_Vi.pptx"
        first.write_text("ORIGINAL CONTENT")

        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "doc_Vi(1).pptx"
        # Original file untouched
        assert first.read_text() == "ORIGINAL CONTENT"

    def test_works_for_xlsx(self, tmp_path: Path):
        src = tmp_path / "book.xlsx"
        src.touch()
        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "book_Vi.xlsx"
        assert out.suffix == ".xlsx"

    def test_works_for_docx(self, tmp_path: Path):
        src = tmp_path / "letter.docx"
        src.touch()
        out = get_unique_output_path(src, "Vietnamese")
        assert out.name == "letter_Vi.docx"
        assert out.suffix == ".docx"

    def test_works_for_pdf(self, tmp_path: Path):
        src = tmp_path / "scan.pdf"
        src.touch()
        out = get_unique_output_path(src, "English", extension=".pdf")
        assert out.name == "scan_En.pdf"

    def test_custom_output_dir_collision(self, tmp_path: Path):
        src = tmp_path / "input.pptx"
        src.touch()
        out_dir = tmp_path / "out"
        out_dir.mkdir()
        (out_dir / "input_Vi.pptx").touch()

        out = get_unique_output_path(src, "Vietnamese", output_dir=out_dir)
        assert out.parent == out_dir
        assert out.name == "input_Vi(1).pptx"


# ═══════════════════════════════════════════════════════════════════════
# Integration: PPT processor uses unique path
# ═══════════════════════════════════════════════════════════════════════

class TestPptProcessorCollision:
    """Verify PptProcessor.process() never overwrites existing output."""

    def _make_pptx(self, path: Path, texts: list[str]) -> Path:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        for i, t in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(0.5))
            txBox.text_frame.text = t
        prs.save(str(path))
        return path

    def test_translate_twice_creates_two_files(self, tmp_path: Path):
        """Translating the same source file twice must create two output files."""
        from app.core.ppt_processor import PptProcessor
        from app.core.translators.translator_service import MockTranslator

        src = self._make_pptx(tmp_path / "input.pptx", ["Hello world"])

        # First translation
        proc1 = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
        )
        out1 = proc1.process(src)
        assert out1.exists()
        assert out1.name == "input_Vi.pptx"

        # Second translation – must NOT overwrite out1
        proc2 = PptProcessor(
            translator=MockTranslator(),
            source_lang="en",
            target_lang="Vietnamese",
        )
        out2 = proc2.process(src)
        assert out2.exists()
        assert out2.name == "input_Vi(1).pptx"

        # Both files still exist
        assert out1.exists()
        assert out2.exists()

    def test_translate_three_times(self, tmp_path: Path):
        """Three translations → _Vi, _Vi(1), _Vi(2)."""
        from app.core.ppt_processor import PptProcessor
        from app.core.translators.translator_service import MockTranslator

        src = self._make_pptx(tmp_path / "slides.pptx", ["Test"])

        outputs = []
        for _ in range(3):
            proc = PptProcessor(
                translator=MockTranslator(),
                source_lang="en",
                target_lang="Vietnamese",
            )
            outputs.append(proc.process(src))

        assert outputs[0].name == "slides_Vi.pptx"
        assert outputs[1].name == "slides_Vi(1).pptx"
        assert outputs[2].name == "slides_Vi(2).pptx"

        # All three files exist
        for o in outputs:
            assert o.exists()

    def test_english_collision(self, tmp_path: Path):
        """Same logic works for _En suffix."""
        from app.core.ppt_processor import PptProcessor
        from app.core.translators.translator_service import MockTranslator

        src = self._make_pptx(tmp_path / "doc.pptx", ["Hello"])

        out1 = PptProcessor(
            translator=MockTranslator(), source_lang="vi", target_lang="English",
        ).process(src)
        out2 = PptProcessor(
            translator=MockTranslator(), source_lang="vi", target_lang="English",
        ).process(src)

        assert out1.name == "doc_En.pptx"
        assert out2.name == "doc_En(1).pptx"
        assert out1.exists()
        assert out2.exists()
