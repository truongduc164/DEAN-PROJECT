"""
Regression tests for custom translated format behavior in Word/Excel.
"""
from __future__ import annotations

import sys
from datetime import datetime
from pathlib import Path

import openpyxl
from docx import Document
from openpyxl.cell.rich_text import CellRichText, TextBlock
from openpyxl.styles import Font
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from app.core.excel_processor import ExcelProcessor
from app.core.pinyin_helper import render_text_with_pinyin
from app.core.ppt_processor import PptProcessor
from app.core.translators.translator_service import MockTranslator
from app.core.word_processor import WordProcessor
from app.settings.settings_manager import settings
import app.core.excel_processor as excel_processor_module
import app.core.ppt_processor as ppt_processor_module
import app.core.word_processor as word_processor_module


def _mock_custom_format_settings(monkeypatch):
    cfg = {
        "text_style_settings.translated_text_format_mode": "custom_format",
        "text_style_settings.keep_format": False,
        "text_style_settings.font_family": "Arial",
        "text_style_settings.font_size": 16,
        "text_style_settings.font_color": "#FF0000",
        "text_style_settings.bold": True,
        "text_style_settings.italic": False,
        "text_style_settings.underline": False,
    }
    monkeypatch.setattr(settings, "get", lambda key, default=None: cfg.get(key, default))


def _mock_pinyin_settings(monkeypatch):
    cfg = {
        "text_style_settings.translated_text_format_mode": "keep_original_format",
        "text_style_settings.keep_format": True,
        "processing_options.add_chinese_pinyin": True,
    }
    monkeypatch.setattr(settings, "get", lambda key, default=None: cfg.get(key, default))


def _run_for_text(paragraph, text: str):
    for run in paragraph.runs:
        if run.text == text:
            return run
    raise AssertionError(f"Run with text '{text}' not found")


def test_word_custom_format_suffix_preserves_original_runs(monkeypatch):
    _mock_custom_format_settings(monkeypatch)

    doc = Document()
    para = doc.add_paragraph()
    r1 = para.add_run("Nguon ")
    r1.font.name = "Calibri"
    r2 = para.add_run("Goc")
    r2.italic = True
    r2.font.name = "Times New Roman"

    proc = WordProcessor(
        translator=MockTranslator(),
        source_lang="English",
        target_lang="Vietnamese",
        output_mode="suffix",
    )
    proc._write_para_result(para, "Nguon Goc", "Ban dich")

    assert para.text == "Nguon Goc\nBan dich"
    assert _run_for_text(para, "Nguon ").font.name == "Calibri"
    assert _run_for_text(para, "Goc").italic is True

    translated_run = _run_for_text(para, "Ban dich")
    assert translated_run.bold is True
    assert translated_run.font.name == "Arial"


def test_word_custom_format_prefix_preserves_original_runs(monkeypatch):
    _mock_custom_format_settings(monkeypatch)

    doc = Document()
    para = doc.add_paragraph()
    r1 = para.add_run("Nguon ")
    r1.bold = True
    r2 = para.add_run("Goc")
    r2.underline = True

    proc = WordProcessor(
        translator=MockTranslator(),
        source_lang="English",
        target_lang="Vietnamese",
        output_mode="prefix",
    )
    proc._write_para_result(para, "Nguon Goc", "Ban dich")

    assert para.text == "Ban dich\nNguon Goc"
    assert _run_for_text(para, "Nguon ").bold is True
    assert _run_for_text(para, "Goc").underline is True

    translated_run = _run_for_text(para, "Ban dich")
    assert translated_run.bold is True
    assert translated_run.font.name == "Arial"


def test_excel_custom_format_suffix_only_styles_translated_part(monkeypatch):
    _mock_custom_format_settings(monkeypatch)

    wb = openpyxl.Workbook()
    ws = wb.active
    cell = ws["A1"]
    cell.value = "Original"
    cell.font = Font(name="Calibri", size=11, bold=False, color="FF0000FF")

    proc = ExcelProcessor(
        translator=MockTranslator(),
        source_lang="English",
        target_lang="Vietnamese",
        output_mode="suffix",
    )
    proc._write_cell_result(cell, "Original", "Translated")

    assert isinstance(cell.value, CellRichText)
    rich = cell.value

    assert isinstance(rich[0], TextBlock)
    assert rich[0].text == "Original"
    assert rich[0].font.rFont == "Calibri"
    assert rich[1] == "\n"
    assert isinstance(rich[2], TextBlock)
    assert rich[2].text == "Translated"
    assert rich[2].font.rFont == "Arial"
    assert rich[2].font.b is True


def test_excel_datetime_cells_are_not_eligible_for_translation():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"].value = datetime(2025, 9, 1, 0, 0, 0)
    ws["A2"].value = "2025-09-01 00:00:00"
    ws["A3"].value = "09/01/2025"
    ws["A4"].value = "13:45:00"
    ws["A5"].value = "Model 2025-AB-01"

    proc = ExcelProcessor(
        translator=MockTranslator(),
        source_lang="English",
        target_lang="Vietnamese",
    )

    assert proc._is_eligible(ws["A1"]) is False
    assert proc._is_eligible(ws["A2"]) is False
    assert proc._is_eligible(ws["A3"]) is False
    assert proc._is_eligible(ws["A4"]) is False
    assert proc._is_eligible(ws["A5"]) is True


def test_ppt_custom_format_suffix_only_styles_translated_part(monkeypatch):
    _mock_custom_format_settings(monkeypatch)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.add_textbox(0, 0, 300, 120)
    para = shape.text_frame.paragraphs[0]

    r1 = para.add_run()
    r1.text = "Nguon "
    r1.font.name = "Calibri"
    r2 = para.add_run()
    r2.text = "Goc"
    r2.font.italic = True
    r2.font.name = "Times New Roman"

    proc = PptProcessor(
        translator=MockTranslator(),
        source_lang="English",
        target_lang="Vietnamese",
        output_mode="suffix",
    )
    proc._write_para_result(para, "Nguon Goc", "Ban dich")

    assert para.text == "Nguon Goc\vBan dich"
    assert _run_for_text(para, "Nguon ").font.name == "Calibri"
    assert _run_for_text(para, "Goc").font.name == "Times New Roman"
    assert _run_for_text(para, "Goc").font.italic is True

    translated_run = _run_for_text(para, "Ban dich")
    assert translated_run.font.name == "Arial"
    assert translated_run.font.bold is True


def test_ppt_custom_format_prefix_only_styles_translated_part(monkeypatch):
    _mock_custom_format_settings(monkeypatch)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.add_textbox(0, 0, 300, 120)
    para = shape.text_frame.paragraphs[0]

    r1 = para.add_run()
    r1.text = "Nguon "
    r1.font.bold = True
    r2 = para.add_run()
    r2.text = "Goc"
    r2.font.underline = True

    proc = PptProcessor(
        translator=MockTranslator(),
        source_lang="English",
        target_lang="Vietnamese",
        output_mode="prefix",
    )
    proc._write_para_result(para, "Nguon Goc", "Ban dich")

    assert para.text == "Ban dich\vNguon Goc"
    assert _run_for_text(para, "Nguon ").font.bold is True
    assert _run_for_text(para, "Goc").font.underline is True

    translated_run = _run_for_text(para, "Ban dich")
    assert translated_run.font.name == "Arial"
    assert translated_run.font.bold is True


def test_excel_suffix_adds_pinyin_under_chinese_source(monkeypatch):
    _mock_pinyin_settings(monkeypatch)
    monkeypatch.setattr(
        excel_processor_module,
        "render_text_with_pinyin",
        lambda text: "你好\nni3 hao3" if text == "你好" else text,
    )

    wb = openpyxl.Workbook()
    ws = wb.active
    cell = ws["A1"]

    proc = ExcelProcessor(
        translator=MockTranslator(),
        source_lang="Chinese",
        target_lang="Vietnamese",
        output_mode="suffix",
    )
    proc._write_cell_result(cell, "你好", "Xin chao")

    assert cell.value == "你好\nni3 hao3\nXin chao"


def test_render_text_with_pinyin_skips_non_chinese_lines(monkeypatch):
    import app.core.pinyin_helper as pinyin_helper_module

    monkeypatch.setattr(
        pinyin_helper_module,
        "get_pinyin_line",
        lambda text: "mei3 tian1" if text == "每天" else "",
    )

    assert render_text_with_pinyin("每天\nNăng suất") == "每天\nmei3 tian1\nNăng suất"


def test_excel_suffix_keeps_pinyin_under_chinese_line_only(monkeypatch):
    _mock_pinyin_settings(monkeypatch)
    monkeypatch.setattr(
        excel_processor_module,
        "get_pinyin_line",
        lambda text: "mei3 tian1" if text == "每天" else "",
    )
    monkeypatch.setattr(
        excel_processor_module,
        "render_text_with_pinyin",
        lambda text: "每天\nmei3 tian1\nNăng suất" if text == "每天\nNăng suất" else text,
    )

    wb = openpyxl.Workbook()
    ws = wb.active
    cell = ws["A1"]

    proc = ExcelProcessor(
        translator=MockTranslator(),
        source_lang="Chinese",
        target_lang="English",
        output_mode="suffix",
    )
    proc._write_cell_result(cell, "每天\nNăng suất", "Daily output")

    assert cell.value == "每天\nmei3 tian1\nNăng suất\nDaily output"


def test_word_suffix_adds_pinyin_under_chinese_source(monkeypatch):
    _mock_pinyin_settings(monkeypatch)
    monkeypatch.setattr(
        word_processor_module,
        "get_pinyin_line",
        lambda text: "ni3 hao3" if text == "你好" else "",
    )

    doc = Document()
    para = doc.add_paragraph()
    para.add_run("你好")

    proc = WordProcessor(
        translator=MockTranslator(),
        source_lang="Chinese",
        target_lang="Vietnamese",
        output_mode="suffix",
    )
    proc._write_para_result(para, "你好", "Xin chao")

    assert para.text == "你好\nni3 hao3\nXin chao"


def test_ppt_suffix_adds_pinyin_under_chinese_source(monkeypatch):
    _mock_pinyin_settings(monkeypatch)
    monkeypatch.setattr(
        ppt_processor_module,
        "get_pinyin_line",
        lambda text: "ni3 hao3" if text == "你好" else "",
    )

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.add_textbox(0, 0, 300, 120)
    para = shape.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "你好"

    proc = PptProcessor(
        translator=MockTranslator(),
        source_lang="Chinese",
        target_lang="Vietnamese",
        output_mode="suffix",
    )
    proc._write_para_result(para, "你好", "Xin chao")

    assert para.text == "你好\vni3 hao3\vXin chao"
