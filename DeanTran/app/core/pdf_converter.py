from __future__ import annotations
import sys
from pathlib import Path
import logging

logger = logging.getLogger("DeanTran.pdf_converter")

def pdf_to_word(
    pdf_path: str | Path, 
    docx_path: str | Path, 
    log_fn=None,
    layout_mode: int = 0,
    parse_header_footer: bool = True,
    parse_table: bool = True,
    parse_image: bool = True
):
    """Convert PDF file to Word (.docx) using pdf2docx with layout options."""
    from pdf2docx import Converter
    
    msg = f"Chuyển đổi PDF sang Word: {Path(pdf_path).name} (Layout: {layout_mode}, Header/Footer: {parse_header_footer}, Table: {parse_table}, Image: {parse_image})"
    logger.info(msg)
    if log_fn:
        log_fn("INFO", msg)
        
    cv = Converter(str(pdf_path))
    cv.convert(
        str(docx_path), 
        start=0, 
        end=None, 
        layout_mode=layout_mode,
        parse_header_footer=parse_header_footer,
        parse_table=parse_table,
        parse_image=parse_image
    )
    cv.close()
    
    msg_done = f"Đã chuyển đổi thành công sang Word: {Path(docx_path).name}"
    logger.info(msg_done)
    if log_fn:
        log_fn("INFO", msg_done)

def pdf_to_excel(
    pdf_path: str | Path, 
    xlsx_path: str | Path, 
    log_fn=None,
    merge_sheets: bool = True,
    auto_fit_columns: bool = True
):
    """Convert PDF file to Excel (.xlsx) extracting both text and tables vertically in order."""
    import pdfplumber
    from openpyxl import Workbook
    from openpyxl.styles import Border, Side
    
    msg = f"Chuyển đổi PDF sang Excel: {Path(pdf_path).name} (Gộp trang: {merge_sheets})"
    logger.info(msg)
    if log_fn:
        log_fn("INFO", msg)
        
    wb = Workbook()
    default_sheet = wb.active
    if default_sheet is not None:
        wb.remove(default_sheet)
        
    thin_border = Border(
        left=Side(style='thin', color='D0D0D0'),
        right=Side(style='thin', color='D0D0D0'),
        top=Side(style='thin', color='D0D0D0'),
        bottom=Side(style='thin', color='D0D0D0')
    )
    
    sheet = None
    if merge_sheets:
        sheet = wb.create_sheet(title="Merged Document")
        
    row_idx = 1
    
    with pdfplumber.open(pdf_path) as pdf:
        for idx, page in enumerate(pdf.pages, 1):
            if not merge_sheets:
                sheet = wb.create_sheet(title=f"Page {idx}")
                row_idx = 1
            else:
                if idx > 1:
                    row_idx += 2 # Thêm dòng trống giữa các trang
                    
            # 1. Tìm tất cả các bảng và bounding box của chúng
            tables = page.find_tables()
            table_bboxes = [t.bbox for t in tables]
            
            # Sắp xếp các bảng theo tọa độ 'top' tăng dần
            sorted_tables = sorted(tables, key=lambda t: t.bbox[1])
            
            # 2. Trích xuất các từ và lọc ra những từ nằm trong các bảng
            words = page.extract_words()
            non_table_words = []
            for w in words:
                in_table = False
                cx = (w['x0'] + w['x1']) / 2.0
                cy = (w['top'] + w['bottom']) / 2.0
                for bbox in table_bboxes:
                    tx0, ttop, tx1, tbottom = bbox
                    if (tx0 <= cx <= tx1) and (ttop <= cy <= tbottom):
                        in_table = True
                        break
                if not in_table:
                    non_table_words.append(w)
            
            # Nhóm các từ ngoài bảng thành các dòng dựa trên tọa độ 'top' (sai số nhỏ, ví dụ 3px)
            lines = []
            if non_table_words:
                non_table_words_sorted = sorted(non_table_words, key=lambda w: w['top'])
                current_line = [non_table_words_sorted[0]]
                for w in non_table_words_sorted[1:]:
                    if w['top'] - current_line[-1]['top'] <= 3.0:
                        current_line.append(w)
                    else:
                        lines.append(current_line)
                        current_line = [w]
                lines.append(current_line)
            
            # Xây dựng các đối tượng dòng văn bản kèm tọa độ 'top'
            text_elements = []
            for line in lines:
                sorted_line = sorted(line, key=lambda w: w['x0'])
                text = " ".join([w['text'] for w in sorted_line])
                top_y = min([w['top'] for w in sorted_line])
                text_elements.append({
                    "type": "text",
                    "top": top_y,
                    "content": text
                })
                
            # Xây dựng các đối tượng bảng kèm tọa độ 'top'
            table_elements = []
            for t in sorted_tables:
                table_elements.append({
                    "type": "table",
                    "top": t.bbox[1],
                    "content": t.extract()
                })
                
            # Gộp và sắp xếp tất cả các phần tử (chữ và bảng) theo tọa độ 'top' từ trên xuống dưới
            all_elements = sorted(text_elements + table_elements, key=lambda e: e['top'])
            
            # 3. Ghi vào Sheet
            for elem in all_elements:
                if elem["type"] == "text":
                    sheet.cell(row=row_idx, column=1, value=elem["content"])
                    row_idx += 1
                elif elem["type"] == "table":
                    table_data = elem["content"]
                    if not table_data:
                        continue
                    
                    for r_data in table_data:
                        for col_idx, val in enumerate(r_data, 1):
                            cell = sheet.cell(row=row_idx, column=col_idx, value=val)
                            cell.border = thin_border
                        row_idx += 1
                    row_idx += 1
                    
            # 4. Tự động giãn cột (Auto-fit Columns)
            if auto_fit_columns and sheet is not None:
                for col in sheet.columns:
                    max_len = 0
                    col_letter = col[0].column_letter
                    for cell in col:
                        if cell.value:
                            val_str = str(cell.value)
                            max_len = max(max_len, len(val_str))
                    if max_len > 0:
                        sheet.column_dimensions[col_letter].width = min(max(max_len + 3, 10), 50)
                        
    wb.save(xlsx_path)
    
    msg_done = f"Đã chuyển đổi thành công sang Excel: {Path(xlsx_path).name}"
    logger.info(msg_done)
    if log_fn:
        log_fn("INFO", msg_done)

def pdf_to_ppt(pdf_path: str | Path, pptx_path: str | Path, log_fn=None, mode: str = "image"):
    """Convert PDF file to PowerPoint (.pptx) rendering as images or parsing as editable text boxes."""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    
    msg = f"Chuyển đổi PDF sang PowerPoint: {Path(pdf_path).name} (Chế độ: {mode})"
    logger.info(msg)
    if log_fn:
        log_fn("INFO", msg)
        
    prs = Presentation()
    doc = fitz.open(pdf_path)
    
    for page_idx in range(len(doc)):
        msg_page = f"Đang dựng slide cho trang {page_idx + 1}/{len(doc)}..."
        logger.info(msg_page)
        if log_fn:
            log_fn("INFO", msg_page)
            
        page = doc[page_idx]
        width_inch = page.rect.width / 72.0
        height_inch = page.rect.height / 72.0
        
        if page_idx == 0:
            prs.slide_width = Inches(width_inch)
            prs.slide_height = Inches(height_inch)
            
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        slide_mode = mode
        if slide_mode == "text":
            # Editable Text Slide mode
            blocks = page.get_text("blocks")
            for b in blocks:
                x0, y0, x1, y1, text, block_no, block_type = b
                if block_type == 0: # Text block
                    text_str = text.strip()
                    if not text_str:
                        continue
                    
                    left = Inches(x0 / 72.0)
                    top = Inches(y0 / 72.0)
                    width = Inches(max(x1 - x0, 10.0) / 72.0)
                    height = Inches(max(y1 - y0, 10.0) / 72.0)
                    
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.text = text_str
                    
            if len(slide.shapes) == 0:
                slide_mode = "image"
                
        if slide_mode == "image":
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            img_stream = BytesIO(img_data)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(width_inch), Inches(height_inch))
            
    prs.save(pptx_path)
    
    msg_done = f"Đã chuyển đổi thành công sang PowerPoint: {Path(pptx_path).name}"
    logger.info(msg_done)
    if log_fn:
        log_fn("INFO", msg_done)
