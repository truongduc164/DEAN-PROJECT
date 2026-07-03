"""
PptProcessor – Translate PowerPoint (.pptx) files.

Extracts text from slides (shapes, text frames, tables, groups),
translates via the same translator pipeline as Excel, and writes
back to a new output file with _Vi/_En suffix.
"""
from __future__ import annotations

import hashlib
import logging
import math
import threading
import time
import traceback
from copy import deepcopy
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.event_manager import EventManager
from app.core.output_naming import (
    get_unique_output_path, build_output_path as _build_base_path,
    get_lang_suffix,
)
from app.core.pinyin_helper import get_pinyin_line, is_available as pinyin_available
from app.term_engine.glossary_loader import GlossaryLoader
from app.term_engine.term_override import TermOverride
from app.term_engine import regex_protection

logger = logging.getLogger("DeanTran.ppt_processor")


# ── Output path (delegates to shared helper) ────────────────────────

def build_ppt_output_path(input_path: Path, target_lang: str) -> Path:
    """Build BASE output path (backward-compat wrapper)."""
    return _build_base_path(input_path, target_lang)


class PptProcessor:
    """Translate a PowerPoint file."""

    def __init__(
        self,
        translator,
        event_manager: Optional[EventManager] = None,
        source_lang: str = "Chinese",
        target_lang: str = "Vietnamese",
        prompt: str = "",
        min_interval: float = 0.2,
        glossary: Optional[GlossaryLoader] = None,
        use_glossary: bool = False,
        use_protection: bool = True,
        translate_textboxes: bool = True,
        pause_event: Optional[threading.Event] = None,
        cancel_event: Optional[threading.Event] = None,
        output_mode: str = "overwrite",
        batch_mode: bool = True,
    ) -> None:
        self.translator = translator
        self.em = event_manager or EventManager()
        self.source_lang = source_lang
        self.target_lang = target_lang
        self.prompt = prompt
        self.min_interval = min_interval
        self.use_protection = use_protection
        self.translate_textboxes = translate_textboxes
        self.output_mode = output_mode

        # Initialize OCR pipeline
        from app.core.ocr.image_translation_pipeline import ImageTranslationPipeline
        self.ocr_pipeline = ImageTranslationPipeline(self.translator, self.em)

        # Auto-disable batch mode for MockTranslator (test safety)
        from app.core.translators.translator_service import MockTranslator
        if isinstance(translator, MockTranslator):
            self.batch_mode = False
        else:
            self.batch_mode = batch_mode

        self._pause_event = pause_event or threading.Event()
        if pause_event is None:
            self._pause_event.set()
        self._cancel_event = cancel_event or threading.Event()

        # Glossary
        self.use_glossary = use_glossary
        self._term_override: Optional[TermOverride] = None

        # Counters
        self._translated_count = 0
        self._error_count = 0
        self._api_calls = 0

        from app.settings.settings_manager import settings

        translated_mode = settings.get("text_style_settings.translated_text_format_mode", None)
        keep_format_legacy = settings.get("text_style_settings.keep_format", True)
        self._text_format_mode = self._resolve_text_format_mode(
            translated_mode,
            keep_format_legacy,
        )
        self._add_chinese_pinyin = settings.get("processing_options.add_chinese_pinyin", False)
        self._pinyin_font_size = settings.get("processing_options.pinyin_font_size", 10)
        self._pinyin_font_family = settings.get("processing_options.pinyin_font_family", "Arial")
        self._pinyin_font_color = settings.get("processing_options.pinyin_font_color", "#888888")
        self._pinyin_format_mode = settings.get("processing_options.pinyin_format_mode", "custom")
        if self._add_chinese_pinyin and not pinyin_available():
            self.em.log("WARN", "Pinyin is enabled but pypinyin is not installed; skipping pinyin.")

    def process(self, file_path: str | Path, output_dir: Optional[str | Path] = None) -> Path:
        file_path = Path(file_path)
        self.em.log("INFO", f"Loading presentation: {file_path.name}")
        prs = Presentation(str(file_path))

        out_path = get_unique_output_path(file_path, self.target_lang, output_dir=output_dir)
        self.em.log("INFO", f"output_path={out_path}")

        # Reset counters
        self._translated_count = 0
        self._error_count = 0
        self._api_calls = 0
        self._last_result = {}
        # Keep track of generated overlay shapes for cleanup before save
        self._created_overlays = []
        t0 = time.time()

        # Collect all text paragraphs & images
        paras_to_translate: list[tuple] = []  # (para, original_text)
        images_to_ocr: list[tuple] = []       # (slide, shape, image_bytes, slide_idx)
        scan_stats = {
            "text_shapes": 0,
            "text_frame_paras": 0,
            "table_paras": 0,
            "images": 0,
            "image_hashes": set(),
        }

        from app.settings.settings_manager import settings
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if self._cancel_event.is_set():
                    break
                self._collect_paras_from_shape(shape, paras_to_translate, scan_stats)
                # Check for images if OCR is globally checked
                if settings.get("ocr_settings.image_text_translation_enabled", False):
                    self._collect_images_from_shape(slide, shape, images_to_ocr, slide_idx, scan_stats)

        total_paras = len(paras_to_translate)
        total_images = len(images_to_ocr)
        unique_images = len(scan_stats["image_hashes"]) if scan_stats["image_hashes"] else 0
        self.em.log("INFO", f"translatable_paras={total_paras}, images={total_images}, unique_images={unique_images}")
        self._log_scan_summary(scan_stats, total_paras, total_images, unique_images)

        if total_paras == 0 and total_images == 0:
            self.em.log("WARN", "No translatable text found!")
            prs.save(str(out_path))
            self.last_result = {
                "status": "SUCCESS",
                "output_path": str(out_path),
                "translated": 0,
                "failed": 0,
                "api_calls": 0,
                "cache_hits": 0,
                "elapsed": time.time() - t0,
            }
            return out_path

        if total_paras > 0:
            # Translate paragraphs
            if self.batch_mode:
                self.em.log("INFO", "Starting BATCH translation mode for PPT")
                self._process_batch_mode(paras_to_translate, total_paras)
            else:
                self.em.log("INFO", "Starting paragraph-by-paragraph fallback mode")
                self._process_cell_by_cell(paras_to_translate, total_paras)

        if total_images > 0:
            self._process_images(images_to_ocr)

        elapsed = time.time() - t0

        # Determine status
        cancelled = self._cancel_event.is_set()
        if cancelled:
            result_status = "CANCELLED"
            stop_name = file_path.stem + "_stop" + file_path.suffix
            out_path = file_path.parent / stop_name
            self.em.log("WARN", f"Cancelled → saving as {stop_name}")
        elif self._translated_count == 0 and self._error_count > 0:
            result_status = "FAILED"
        elif self._error_count > 0:
            result_status = "PARTIAL"
        else:
            result_status = "SUCCESS"

        if getattr(self, "_created_overlays", None):
            self._cleanup_ocr_overlays()

        # Save
        out_path = self._safe_save(prs, out_path)

        self.em.log(
            "INFO",
            f"result={result_status} translated={self._translated_count} "
            f"errors={self._error_count} elapsed={elapsed:.1f}s "
            f"output_path={out_path}",
        )

        self.last_result = {
            "status": result_status,
            "output_path": str(out_path),
            "translated": self._translated_count,
            "failed": self._error_count,
            "api_calls": self._api_calls,
            "cache_hits": 0,
            "elapsed": elapsed,
        }

        return out_path

    def _process_batch_mode(self, paras_to_translate: list[tuple], total: int):
        from app.core.key_provider import KeyProvider
        from app.core.translators.batch_translator import (
            GeminiBatchTranslator, CellItem,
        )
        from app.settings.settings_manager import settings

        model = settings.get("selected_models.gemini", "gemini-3.1-flash-lite")
        kp = KeyProvider()
        kp.log_status()

        # Check if DeepSeek key is available (for DeepSeek tool)
        _tool = settings.get("translation_tool", "Gemini")
        _has_deepseek_key = False
        if _tool == "DeepSeek":
            try:
                import os
                from app.core.secure_storage import SecureStorage
                _ds_key = SecureStorage().load_deepseek_key() or os.environ.get("DEEPSEEK_API_KEY", "")
                _has_deepseek_key = bool(_ds_key and _ds_key.strip())
            except Exception:
                pass
        _keys_available = kp.key_loaded or _has_deepseek_key

        # Prepare CellItems
        items: list[CellItem] = []
        paras_map = {}
        for idx, (para, text) in enumerate(paras_to_translate):
            cell_id = f"para_{idx}"
            paras_map[cell_id] = para
            
            item = CellItem(cell_id=cell_id, original=text)
            item.real_original = text

            if self.use_glossary and getattr(self, '_term_override', None):
                term_dict = self._term_override.glossary.get_dict() if self._term_override.glossary else {}
                matched_val = None
                
                lower_text = text.strip().lower()
                for k, v in term_dict.items():
                    if k.lower() == lower_text:
                        matched_val = v
                        break
                        
                replaced_text = self._term_override.pre_replace(text)
                
                if matched_val:
                    item.translated = matched_val
                elif replaced_text != text and not _keys_available:
                    item.translated = replaced_text
                    
                item.original = replaced_text
                
            if not _keys_available and not item.translated:
                item.translated = item.original
                item.error = "No API key available"

            items.append(item)
            
        if not _keys_available:
            self.em.log("WARN", "No API key available. Only dictionary terms will be translated.")

        bt = GeminiBatchTranslator(
            key_provider=kp,
            model_name=model,
            source_lang=self.source_lang,
            target_lang=self.target_lang,
            prompt=self.prompt,
            min_interval=self.min_interval,
            log_fn=lambda lvl, msg: self.em.log(lvl, msg),
            pause_event=self._pause_event,
            cancel_event=self._cancel_event,
            progress_fn=lambda cur, tot: self.em.progress(cur, tot),
        )

        result = bt.translate_batch(items)
        self._api_calls = result.api_calls

        self.em.log(
            "INFO",
            f"Batch complete: api_calls={result.api_calls} "
            f"cache_hits={result.cache_hits} errors={result.errors} "
            f"elapsed={result.elapsed:.1f}s",
        )

        # Write results back
        done = 0
        for item in items:
            para = paras_map.get(item.cell_id)
            if not para:
                continue

            if item.translated:
                translated = item.translated
                if self.use_glossary and getattr(self, '_term_override', None):
                    translated = self._term_override.post_fix(translated)
                real_orig = getattr(item, 'real_original', item.original)
                if translated.strip() != real_orig.strip():
                    self._write_para_result(para, real_orig, translated)
                self._translated_count += 1
            else:
                self._error_count += 1

            done += 1
            self.em.progress(done, total)

        if result.errors > 0 and self._error_count == 0:
            self._error_count += result.errors

    def _process_cell_by_cell(self, paras_to_translate: list[tuple], total: int):
        done = 0
        for para, original in paras_to_translate:
            if self._cancel_event.is_set():
                break

            self._pause_event.wait()

            try:
                translated = self._translate_text(original)
                if translated.strip() != original.strip():
                    self._write_para_result(para, original, translated)
                self._translated_count += 1
            except Exception as exc:
                self.em.log(
                    "ERROR",
                    f"Translation failed: {type(exc).__name__}: {exc}",
                )
                self._error_count += 1

            done += 1
            self.em.progress(done, total)

            if self.min_interval > 0 and done < total:
                time.sleep(self.min_interval)

    def _collect_paras_from_shape(self, shape, paras_list, scan_stats=None):
        """Recursively collect text paragraphs from a shape."""
        # Group shapes
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                self._collect_paras_from_shape(child, paras_list, scan_stats)
            return

        # Text boxes / auto shapes with text
        if getattr(shape, "has_text_frame", False):
            if not self.translate_textboxes:
                # Skip text boxes (non-placeholder) when user disabled textbox translation
                # Placeholders (title, subtitle, body) are still translated
                try:
                    _ = shape.placeholder_format
                    is_placeholder = True
                except (ValueError, AttributeError):
                    is_placeholder = False
                if not is_placeholder:
                    return

            added_count = 0
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and not text.replace(" ", "").isnumeric():
                    paras_list.append((para, para.text))
                    added_count += 1

            if scan_stats is not None and added_count > 0:
                scan_stats["text_shapes"] += 1
                scan_stats["text_frame_paras"] += added_count

        # Tables
        if getattr(shape, "has_table", False):
            table_added = 0
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and not text.replace(" ", "").isnumeric():
                            paras_list.append((para, para.text))
                            table_added += 1
            if scan_stats is not None and table_added > 0:
                scan_stats["table_paras"] += table_added

    def _collect_images_from_shape(self, slide, shape, images_list, slide_idx=0, scan_stats=None):
        """Recursively collect images from a shape."""
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                self._collect_images_from_shape(slide, child, images_list, slide_idx, scan_stats)
            return

        if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
            try:
                image_bytes = shape.image.blob
                images_list.append((slide, shape, image_bytes, slide_idx))
                if scan_stats is not None:
                    scan_stats["images"] += 1
                    scan_stats["image_hashes"].add(hashlib.sha1(image_bytes).hexdigest())
            except Exception:
                pass

    def _log_scan_summary(self, scan_stats: dict, total_paras: int, total_images: int, unique_images: int) -> None:
        from app.settings.settings_manager import settings

        ocr_enabled = settings.get("ocr_settings.image_text_translation_enabled", False)
        engine = settings.get("ocr_settings.engine", "paddle")
        dedupe_enabled = settings.get("ocr_settings.google_vision_dedupe_images", True)
        batch_enabled = settings.get("ocr_settings.google_vision_batch_enabled", True)
        canvas_enabled = settings.get("ocr_settings.google_vision_canvas_enabled", False)

        ocr_mode = "off"
        est_vision_calls = 0

        considered_images = unique_images if dedupe_enabled else total_images
        if ocr_enabled and engine == "google_vision" and considered_images > 0:
            if canvas_enabled:
                per_canvas = max(1, int(settings.get("ocr_settings.google_vision_canvas_images_per_canvas", 4)))
                est_vision_calls = math.ceil(considered_images / per_canvas)
                ocr_mode = "canvas"
            elif batch_enabled:
                per_request = max(1, int(settings.get("ocr_settings.google_vision_max_images_per_request", 16)))
                est_vision_calls = math.ceil(considered_images / per_request)
                ocr_mode = "batch"
            else:
                est_vision_calls = considered_images
                ocr_mode = "single"
        elif ocr_enabled and engine != "google_vision":
            ocr_mode = engine

        self.em.log(
            "INFO",
            "ppt_scan_summary: "
            f"text_shapes={scan_stats.get('text_shapes', 0)} "
            f"text_frame_paras={scan_stats.get('text_frame_paras', 0)} "
            f"table_paras={scan_stats.get('table_paras', 0)} "
            f"translatable_paras={total_paras} "
            f"images={total_images} unique_images={unique_images} "
            f"ocr_mode={ocr_mode} est_vision_calls={est_vision_calls}",
        )

    @staticmethod
    def _rects_overlap(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        return (ax < bx + bw) and (ax + aw > bx) and (ay < by + bh) and (ay + ah > by)

    @staticmethod
    def _overlap_area(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> int:
        if not PptProcessor._rects_overlap(a, b):
            return 0
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        x_overlap = max(0, min(ax + aw, bx + bw) - max(ax, bx))
        y_overlap = max(0, min(ay + ah, by + bh) - max(ay, by))
        return x_overlap * y_overlap

    def _overlap_score(self, rect: tuple[int, int, int, int], occupied_rects: list[tuple[int, int, int, int]]) -> int:
        return sum(self._overlap_area(rect, occ) for occ in occupied_rects)

    @staticmethod
    def _rect_contains(
        container: tuple[int, int, int, int],
        rect: tuple[int, int, int, int],
    ) -> bool:
        cx, cy, cw, ch = container
        rx, ry, rw, rh = rect
        return rx >= cx and ry >= cy and (rx + rw) <= (cx + cw) and (ry + rh) <= (cy + ch)

    @staticmethod
    def _rect_contains_point(
        container: tuple[int, int, int, int],
        px: int,
        py: int,
    ) -> bool:
        cx, cy, cw, ch = container
        return cx <= px <= (cx + cw) and cy <= py <= (cy + ch)

    def _choose_smart_position(
        self,
        base_rect: tuple[int, int, int, int],
        occupied_rects: list[tuple[int, int, int, int]],
        step_x: int,
        step_y: int,
        max_shift_x: int,
        max_shift_y: int,
    ) -> tuple[tuple[int, int, int, int], bool]:
        base_score = self._overlap_score(base_rect, occupied_rects)
        best_rect = base_rect
        best_score = base_score
        if base_score == 0:
            return base_rect, False

        self.em.log("INFO", "overlap -> shift right")
        for shift in range(step_x, max_shift_x + step_x, step_x):
            cand = (base_rect[0] + shift, base_rect[1], base_rect[2], base_rect[3])
            score = self._overlap_score(cand, occupied_rects)
            if score < best_score:
                best_rect, best_score = cand, score
            if score == 0:
                return cand, False

        self.em.log("INFO", "overlap -> shift down")
        for shift in range(step_y, max_shift_y + step_y, step_y):
            cand = (base_rect[0], base_rect[1] + shift, base_rect[2], base_rect[3])
            score = self._overlap_score(cand, occupied_rects)
            if score < best_score:
                best_rect, best_score = cand, score
            if score == 0:
                return cand, False

        return best_rect, best_score > 0

    def _choose_whitespace_position(
        self,
        base_rect: tuple[int, int, int, int],
        occupied_rects: list[tuple[int, int, int, int]],
        step_x: int,
        step_y: int,
        max_search_x: int,
        max_search_y: int,
        constrain_region: tuple[int, int, int, int] | None = None,
    ) -> tuple[tuple[int, int, int, int] | None, str | None]:
        x, y, w, h = base_rect
        directions = [
            ("right", step_x, max_search_x, lambda s: (x + s, y, w, h)),
            ("down", step_y, max_search_y, lambda s: (x, y + s, w, h)),
            ("left", step_x, max_search_x, lambda s: (x - s, y, w, h)),
            ("up", step_y, max_search_y, lambda s: (x, y - s, w, h)),
        ]

        for name, step, max_dist, make_rect in directions:
            for shift in range(step, max_dist + step, step):
                cand = make_rect(shift)
                if constrain_region is not None and not self._rect_contains(constrain_region, cand):
                    continue
                if self._overlap_score(cand, occupied_rects) == 0:
                    return cand, name
        return None, None

    def _find_same_section_region(
        self,
        base_rect: tuple[int, int, int, int],
        source_rect: tuple[int, int, int, int],
        shape_rects: list[tuple[int, int, int, int]],
    ) -> tuple[int, int, int, int] | None:
        bx, by, bw, bh = base_rect
        center_x = bx + (bw // 2)
        center_y = by + (bh // 2)
        base_area = max(1, bw * bh)

        candidates: list[tuple[int, tuple[int, int, int, int]]] = []
        for rect in shape_rects:
            if rect == source_rect:
                continue
            _, _, rw, rh = rect
            area = rw * rh
            if area < int(base_area * 1.5):
                continue
            if not self._rect_contains_point(rect, center_x, center_y):
                continue
            if not self._rect_contains(rect, base_rect):
                continue
            candidates.append((area, rect))

        if not candidates:
            return None
        candidates.sort(key=lambda item: item[0])
        return candidates[0][1]

    def _resolve_placement_rect(
        self,
        base_rect: tuple[int, int, int, int],
        occupied_rects: list[tuple[int, int, int, int]],
        shape_rects: list[tuple[int, int, int, int]],
        source_rect: tuple[int, int, int, int],
        placement_mode: str,
        step_x: int,
        step_y: int,
        max_shift_x: int,
        max_shift_y: int,
        whitespace_search_x: int,
        whitespace_search_y: int,
    ) -> tuple[tuple[int, int, int, int], bool]:
        if placement_mode == "exact":
            return base_rect, False

        if placement_mode == "whitespace":
            same_section = self._find_same_section_region(base_rect, source_rect, shape_rects)
            if same_section is not None:
                section_occupied = [r for r in occupied_rects if r != same_section]
                same_rect, _same_dir = self._choose_whitespace_position(
                    base_rect,
                    section_occupied,
                    step_x,
                    step_y,
                    whitespace_search_x,
                    whitespace_search_y,
                    constrain_region=same_section,
                )
                if same_rect is not None:
                    self.em.log("INFO", "placement=same_section_whitespace")
                    return same_rect, False

            ws_rect, ws_dir = self._choose_whitespace_position(
                base_rect, occupied_rects, step_x, step_y, whitespace_search_x, whitespace_search_y
            )
            if ws_rect is not None and ws_dir is not None:
                self.em.log("INFO", "placement=nearby_whitespace")
                return ws_rect, False

            self.em.log("INFO", "fallback=smart_adjust")
            smart_rect, still_overlap = self._choose_smart_position(
                base_rect, occupied_rects, step_x, step_y, max_shift_x, max_shift_y
            )
            if still_overlap:
                self.em.log("INFO", "fallback=exact")
                return base_rect, False
            return smart_rect, False

        if placement_mode == "smart_adjust":
            smart_rect, still_overlap = self._choose_smart_position(
                base_rect, occupied_rects, step_x, step_y, max_shift_x, max_shift_y
            )
            if still_overlap:
                self.em.log("INFO", "fallback=exact")
                return base_rect, False
            return smart_rect, False

        return self._choose_smart_position(
            base_rect, occupied_rects, step_x, step_y, max_shift_x, max_shift_y
        )

    @staticmethod
    def _resolve_text_format_mode(
        translated_text_format_mode: str | None,
        keep_format_legacy: bool,
    ) -> str:
        if translated_text_format_mode in ("keep_original_format", "custom_format"):
            return translated_text_format_mode
        return "keep_original_format" if keep_format_legacy else "custom_format"

    @staticmethod
    def _resolve_ocr_placement_mode(
        placement_mode: str | None,
        legacy_mode: str | None,
    ) -> str:
        if placement_mode in ("exact", "smart_adjust", "whitespace"):
            return placement_mode
        if legacy_mode in ("exact", "smart_adjust"):
            return legacy_mode
        return "whitespace"

    @staticmethod
    def _clone_run_props(run):
        rpr = getattr(run._r, "rPr", None)
        return deepcopy(rpr) if rpr is not None else None

    @staticmethod
    def _apply_run_props(run, rpr) -> None:
        if rpr is None:
            return
        existing_rpr = getattr(run._r, "rPr", None)
        if existing_rpr is not None:
            run._r.remove(existing_rpr)
        run._r.insert(0, deepcopy(rpr))

    @staticmethod
    def _snapshot_original_content(para, original: str) -> list[tuple[str, str, object]]:
        snapshot: list[tuple[str, str, object]] = []
        for child in para._p.content_children:
            local_name = child.tag.rsplit("}", 1)[-1]
            if local_name == "br":
                snapshot.append(("break", "", None))
                continue

            text = child.text or ""
            if not text:
                continue

            rpr = getattr(child, "rPr", None)
            snapshot.append(("text", text, deepcopy(rpr) if rpr is not None else None))

        if snapshot or not original:
            return snapshot

        normalized = original.replace("\r", "").replace("\n", "\v")
        for idx, chunk in enumerate(normalized.split("\v")):
            if idx > 0:
                snapshot.append(("break", "", None))
            if chunk:
                snapshot.append(("text", chunk, None))
        return snapshot

    def _append_snapshot_content(self, para, snapshot: list[tuple[str, str, object]]) -> None:
        for kind, text, rpr in snapshot:
            if kind == "break":
                para.add_line_break()
                continue

            run = para.add_run()
            run.text = text
            self._apply_run_props(run, rpr)

    def _pinyin_for_lang(self, text: str, lang: str) -> str:
        if not self._add_chinese_pinyin or lang != "Chinese":
            return ""
        return get_pinyin_line(text)

    def _append_styled_text_block(self, para, text: str, pinyin: str, first_rpr) -> None:
        run = para.add_run()
        run.text = text
        self._style_translated_run(run, first_rpr)
        if pinyin:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            para.add_line_break()
            pinyin_run = para.add_run()
            pinyin_run.text = pinyin
            self._style_translated_run(pinyin_run, first_rpr)
            if self._pinyin_format_mode != "keep_original":
                pinyin_run.font.size = Pt(self._pinyin_font_size)
                pinyin_run.font.name = self._pinyin_font_family
                try:
                    c = self._pinyin_font_color.lstrip("#")
                    pinyin_run.font.color.rgb = RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                except Exception:
                    pass

    def _append_original_block(self, para, snapshot: list[tuple[str, str, object]], pinyin: str, first_rpr) -> None:
        self._append_snapshot_content(para, snapshot)
        if pinyin:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            para.add_line_break()
            pinyin_run = para.add_run()
            pinyin_run.text = pinyin
            self._apply_run_props(pinyin_run, first_rpr)
            if self._pinyin_format_mode != "keep_original":
                pinyin_run.font.size = Pt(self._pinyin_font_size)
                pinyin_run.font.name = self._pinyin_font_family
                try:
                    c = self._pinyin_font_color.lstrip("#")
                    pinyin_run.font.color.rgb = RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                except Exception:
                    pass

    def _apply_custom_run_style(self, run) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        from app.core.ocr.image_text_renderer import _hex_to_rgb
        from app.settings.settings_manager import settings

        font_name = settings.get("text_style_settings.font_family", "Arial")
        font_size = settings.get("text_style_settings.font_size", 14)
        bold = settings.get("text_style_settings.bold", False)
        italic = settings.get("text_style_settings.italic", False)
        underline = settings.get("text_style_settings.underline", False)
        font_color = settings.get("text_style_settings.font_color", "#000000")
        if font_name != "Mặc định":
            run.font.name = font_name
        if font_size > 0:
            run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        if font_color != "default":
            r, g, b = _hex_to_rgb(font_color)
            run.font.color.rgb = RGBColor(r, g, b)

    def _style_translated_run(self, run, first_rpr) -> None:
        if self._text_format_mode == "custom_format":
            self._apply_custom_run_style(run)
        else:
            self._apply_run_props(run, first_rpr)

    def _process_images(self, images_to_ocr: list[tuple]):
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from app.settings.settings_manager import settings
        from app.core.ocr.image_text_renderer import _hex_to_rgb

        self.em.log("INFO", f"Starting OCR batch pipeline for {len(images_to_ocr)} images")

        # ── Step 1: OCR all images + filter segments ─────────────────
        ocr_results = self.ocr_pipeline.ocr_and_filter_images(images_to_ocr, self.source_lang)

        # ── Step 2: Batch translate ALL segments via single Gemini JSON ─
        self.ocr_pipeline.batch_translate_segments(ocr_results, self.source_lang, self.target_lang)

        # ── Step 3: Render overlays ──────────────────────────────────
        if not settings.get("ocr_settings.render_textbox_overlay", True):
            self.em.log("INFO", "render_textbox_overlay=False, skipping overlay rendering")
            return

        placement_mode = self._resolve_ocr_placement_mode(
            settings.get("ocr_settings.ocr_textbox_placement_mode", None),
            settings.get("ocr_settings.ocr_textbox_mode", None),
        )
        format_mode = self._resolve_text_format_mode(
            settings.get("text_style_settings.translated_text_format_mode", None),
            settings.get("text_style_settings.keep_format", True),
        )

        self.em.log("INFO", f"text_format_mode={format_mode}")
        self.em.log("INFO", f"placement_mode={placement_mode}")

        use_custom_format = format_mode == "custom_format"
        if use_custom_format:
            font_name = settings.get("text_style_settings.font_family", "Arial")
            base_font_size = settings.get("text_style_settings.font_size", 14)
            bold = settings.get("text_style_settings.bold", False)
            italic = settings.get("text_style_settings.italic", False)
            underline = settings.get("text_style_settings.underline", False)
            font_color = settings.get("text_style_settings.font_color", "#000000")
            self.em.log(
                "INFO",
                f"custom_font_applied family={font_name} size={base_font_size} color={font_color}",
            )
        else:
            # OCR text has no reliable source run-style, so keep a neutral readable fallback.
            font_name = "Arial"
            base_font_size = 14
            bold = False
            italic = False
            underline = False
            font_color = "#000000"
        
        display_mode = settings.get("ocr_settings.ocr_display_mode", "overwrite")
        display_container = settings.get("ocr_settings.ocr_display_container", "shape")

        r_c, g_c, b_c = _hex_to_rgb(font_color)
        bg_r, bg_g, bg_b = 255, 255, 255

        # Collision tracking per slide:
        # OCR regions, images, shapes, and placed translation textboxes.
        slide_occupied: dict[int, list[tuple[int, int, int, int]]] = {}
        slide_shape_rects: dict[int, list[tuple[int, int, int, int]]] = {}
        for item in ocr_results:
            slide = item["slide"]
            shape = item["shape"]
            result = item["result"]
            slide_id = id(slide)

            if slide_id not in slide_occupied:
                occupied = []
                shape_rects = []
                for shp in slide.shapes:
                    try:
                        rect = (int(shp.left), int(shp.top), int(shp.width), int(shp.height))
                        occupied.append(rect)
                        shape_rects.append(rect)
                    except Exception:
                        continue
                slide_occupied[slide_id] = occupied
                slide_shape_rects[slide_id] = shape_rects

            if result.error or not result.blocks or result.width == 0 or result.height == 0:
                continue

            scale_x = shape.width / result.width
            scale_y = shape.height / result.height
            occupied = slide_occupied[slide_id]
            for block in result.blocks:
                x, y, w, h = block.bbox
                left = int(shape.left + (x * scale_x))
                top = int(shape.top + (y * scale_y))
                width = max(1, int(w * scale_x))
                height = max(1, int(h * scale_y))
                occupied.append((left, top, width, height))

        overlay_count = 0

        for item in ocr_results:
            slide = item["slide"]
            shape = item["shape"]
            result = item["result"]

            if result.error or not result.blocks or result.width == 0 or result.height == 0:
                continue

            if self._cancel_event.is_set():
                break

            scale_x = shape.width / result.width
            scale_y = shape.height / result.height
            occupied = slide_occupied.setdefault(id(slide), [])
            shape_rects = slide_shape_rects.setdefault(id(slide), [])
            source_rect = (int(shape.left), int(shape.top), int(shape.width), int(shape.height))

            for block in result.blocks:
                if block.classification == "noise" or getattr(block, "skipped", False):
                    continue

                original_text = block.text.strip() if block.text else ""
                translated_text = block.translated_text.strip() if block.translated_text else ""
                
                if not original_text and not translated_text:
                    continue
                    
                original_display = original_text
                translated_display = translated_text
                original_pinyin = self._pinyin_for_lang(original_text, self.source_lang)
                translated_pinyin = self._pinyin_for_lang(translated_text, self.target_lang)

                if original_pinyin:
                    original_display = f"{original_text}\n{original_pinyin}"
                if translated_pinyin:
                    translated_display = f"{translated_text}\n{translated_pinyin}"

                if display_mode == "prefix":
                    rendered_text = (
                        f"{translated_display}\n{original_display}"
                        if translated_text and translated_text != original_text
                        else original_display
                    )
                elif display_mode == "suffix":
                    rendered_text = (
                        f"{original_display}\n{translated_display}"
                        if translated_text and translated_text != original_text
                        else original_display
                    )
                else: # overwrite
                    rendered_text = translated_display or original_display

                if not rendered_text.strip():
                    continue

                x, y, w, h = block.bbox

                tb_left = int(shape.left + (x * scale_x))
                tb_top = int(shape.top + (y * scale_y))
                tb_width = max(1, int(w * scale_x))
                tb_height_raw = max(1, int(h * scale_y))

                computed_size = base_font_size
                tb_height = max(tb_height_raw, int(Pt(computed_size * 2)))

                placed_rect = (tb_left, tb_top, tb_width, tb_height)
                still_overlapping = False
                step_x = max(1, int(round(10 * scale_x)))
                step_y = max(1, int(round(10 * scale_y)))
                smart_max_shift_px = max(10, int(settings.get("ocr_settings.smart_adjust_max_shift_px", 80)))
                whitespace_search_px = max(10, int(settings.get("ocr_settings.whitespace_search_max_shift_px", 180)))
                max_shift_x = max(step_x, int(round(smart_max_shift_px * scale_x)))
                max_shift_y = max(step_y, int(round(smart_max_shift_px * scale_y)))
                whitespace_search_x = max(step_x, int(round(whitespace_search_px * scale_x)))
                whitespace_search_y = max(step_y, int(round(whitespace_search_px * scale_y)))
                placed_rect, still_overlapping = self._resolve_placement_rect(
                    placed_rect,
                    occupied,
                    shape_rects,
                    source_rect,
                    placement_mode,
                    step_x,
                    step_y,
                    max_shift_x,
                    max_shift_y,
                    whitespace_search_x,
                    whitespace_search_y,
                )

                self.em.log(
                    "INFO",
                    f"placed bbox=({placed_rect[0]},{placed_rect[1]},{placed_rect[2]},{placed_rect[3]})",
                )

                try:
                    if display_container == "shape":
                        txBox = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            placed_rect[0],
                            placed_rect[1],
                            placed_rect[2],
                            placed_rect[3],
                        )
                        # Remove default outline
                        txBox.line.fill.background()
                    else:
                        txBox = slide.shapes.add_textbox(
                            placed_rect[0],
                            placed_rect[1],
                            placed_rect[2],
                            placed_rect[3],
                        )
                        
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.clear()

                    # Background fill for Shape container
                    if display_container == "shape":
                        fill = txBox.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)

                    p = tf.paragraphs[0]

                    run = p.add_run()
                    run.text = rendered_text
                    run.font.size = Pt(computed_size)
                    run.font.color.rgb = RGBColor(r_c, g_c, b_c)
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.underline = underline
                    run.font.name = font_name
                    
                    block.overlay_created = True
                    self._created_overlays.append((slide, txBox, block, rendered_text))
                    occupied.append(placed_rect)
                    overlay_count += 1
                    if still_overlapping:
                        self.em.log(
                            "WARN",
                            f"[{result.image_id}] smart_adjust could not fully avoid overlap; kept best position",
                        )

                except Exception as e:
                    self.em.log("WARN", f"[{result.image_id}] Failed to overlay box: {e}")

        self.em.log("INFO", f"overlay_textboxes_created={overlay_count}")

    def _write_para_result(self, para, original: str, translated: str):
        """Write translation into a paragraph based on output_mode.
        """
        original = (original or "").replace("\r", "")
        translated = (translated or "").replace("\r", "")

        # Deduplicate identical lines in prefix/suffix mode to avoid duplicating untranslated lines
        if self.output_mode in ("prefix", "suffix") and '\n' in original:
            orig_lines = [l.strip().lower() for l in original.split('\n')]
            trans_lines = translated.split('\n')
            new_trans_lines = []
            for line in trans_lines:
                if line.strip().lower() in orig_lines:
                    continue
                new_trans_lines.append(line)
            if new_trans_lines:
                translated = '\n'.join(new_trans_lines)
            else:
                translated = original

        first_rpr = self._clone_run_props(para.runs[0]) if para.runs else None
        original_snapshot = self._snapshot_original_content(para, original)
        original_pinyin = self._pinyin_for_lang(original, self.source_lang)
        translated_pinyin = self._pinyin_for_lang(translated, self.target_lang)

        para.clear()

        if translated == original:
            self._append_styled_text_block(para, translated, translated_pinyin, first_rpr)
            return

        if self.output_mode == "prefix":
            self._append_styled_text_block(para, translated, translated_pinyin, first_rpr)
            if original_snapshot:
                para.add_line_break()
                self._append_original_block(para, original_snapshot, original_pinyin, first_rpr)
            return

        if self.output_mode == "suffix":
            self._append_original_block(para, original_snapshot, original_pinyin, first_rpr)
            if translated:
                if original_snapshot:
                    para.add_line_break()
                self._append_styled_text_block(para, translated, translated_pinyin, first_rpr)
            return

        self._append_styled_text_block(para, translated, translated_pinyin, first_rpr)

    def _translate_text(self, text: str) -> str:
        if self.use_glossary and self._term_override:
            text = self._term_override.pre_replace(text)
        mapping: Dict[str, str] = {}
        if self.use_protection:
            text, mapping = regex_protection.protect(text)

        translated = self.translator.translate(
            text, self.source_lang, self.target_lang,
        )
        self._api_calls += 1

        if self.use_protection and mapping:
            translated = regex_protection.restore(translated, mapping)
        if self.use_glossary and self._term_override:
            translated = self._term_override.post_fix(translated)
        return translated

    def _safe_save(self, prs: Presentation, out_path: Path) -> Path:
        """Save presentation. If PermissionError, try numbered variants."""
        for attempt in range(10):
            target = out_path if attempt == 0 else (
                out_path.parent / f"{out_path.stem}({attempt}){out_path.suffix}"
            )
            try:
                prs.save(str(target))
                if attempt > 0:
                    self.em.log("WARN", f"File locked, saved as '{target.name}'")
                return target
            except PermissionError:
                self.em.log("WARN", f"PermissionError on '{target.name}', trying next…")
        self.em.log("ERROR", f"Cannot save after 10 attempts: {out_path.name}")
        return out_path

    def _cleanup_ocr_overlays(self):
        """Scan tracking list and remove noisy overlays before final save.
        
        Only removes overlays that are genuinely noise (empty text or
        classified as noise by the OCR pipeline). Safe for valid translations.
        """
        from app.settings.settings_manager import settings
        
        if not settings.get("ocr_settings.remove_noise_overlays_before_save", True):
            return

        removed_count = 0
        for slide, shape, block, rendered_text in self._created_overlays:
            should_remove = False
            clean_text = rendered_text.strip()
            
            if block.classification == "noise":
                should_remove = True
            elif not clean_text:
                should_remove = True
            
            if should_remove:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    removed_count += 1
                    block.overlay_created = False
                    self.em.log("INFO", f"OCR overlay removed during cleanup text={clean_text!r}")
                except Exception as e:
                    self.em.log("WARN", f"Failed to remove noisy overlay: {e}")
                    
        self.em.log("INFO", f"overlays_removed_noise={removed_count}")
