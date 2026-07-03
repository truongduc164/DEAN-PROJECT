import sys
from pathlib import Path
from pptx import Presentation

# Find the file dynamically to avoid unicode string hardcoding issues
base_dir = Path(r"D:\0. Lập trình\1.DEANTRANS\outputs")
target_file = None
for f in base_dir.glob("*.pptx"):
    if f.name.startswith("~$"): continue
    if "38" in f.name and "QC" in f.name and not f.name.endswith("_En.pptx") and not f.name.endswith("_Vi.pptx"):
        target_file = f
        break

if not target_file:
    print("File not found via glob.")
    sys.exit(1)

# skip print
prs = Presentation(str(target_file))

with open("ppt_inspect.log", "w", encoding="utf-8") as f:
    for slide_idx, slide in enumerate(prs.slides, 1):
        f.write(f"\n--- Slide {slide_idx} ---\n")
        for shape_idx, shape in enumerate(slide.shapes):
            f.write(f"  Shape {shape_idx}: type={shape.shape_type}, name={shape.name}\n")
            if getattr(shape, "has_text_frame", False):
                f.write(f"    has_text_frame=True\n")
                try:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        f.write(f"      Para {para_idx}: runs={len(para.runs)}\n")
                        for r_idx, r in enumerate(para.runs):
                            # Ensure we handle potentially empty runs or weird encodings safely
                            safe_text = r.text.encode("unicode_escape").decode("utf-8")
                            f.write(f"        Run {r_idx}: text='{safe_text}'\n")
                            # Check if this run contains an inline drawing/picture
                            xml = r._r.xml
                            if "a:blip" in xml or "a:drawing" in xml:
                                f.write(f"          [!] RUN CONTAINS IMAGE/DRAWING!\n")
                except Exception as e:
                    f.write(f"      Error iterating text frame: {e}\n")
            
            # Check for grouped shapes
            if shape.shape_type == 6: # GROUP
                f.write("    (Group Shape)\n")
                for child in shape.shapes:
                    f.write(f"      Child: type={child.shape_type}, name={child.name}\n")
print("Done inspecting.")
